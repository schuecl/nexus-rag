"""Issue #92 (split from #88): caption embedded images/figures at ingestion
time so their content becomes retrievable text instead of silently vanishing
from the corpus.

Design decisions the issue asked to be made explicitly:

*Where it runs* -- inline in the worker's processing pass, between parse and
chunk, not as a separate queue stage. A distinct JetStream stage would need
its own consumer, redelivery semantics, and a second pending state for a
feature that is optional and bounded; inline, the captions simply become more
ParsedSections flowing through the existing chunk/embed/store path, and the
durable-redelivery story stays exactly #164's.

*Failure handling* -- degrade, never fail (the reranker's degrade-on-outage
pattern, orchestration-mcp/app/reranking.py, not ParsingError's fail-the-
document pattern). Text extraction succeeded; discarding it because a
captioning sidecar is down would strictly worsen the corpus. Every failure
path here -- model unavailable, per-image error, the whole pass outrunning its
budget -- logs, counts (metrics.images_skipped_total), and returns whatever
captions were produced so far. This module never raises to its caller.

*Time bounds* -- captioning runs inside #208's PROCESSING_TIMEOUT_SECONDS
budget, so it must not be able to eat it: the whole pass is additionally
bounded by CAPTIONING_TIMEOUT_SECONDS (default 90s, well under the 240s
document budget), each Ollama call by CAPTION_REQUEST_TIMEOUT_SECONDS, and
the image count by MAX_IMAGES_PER_DOCUMENT. Worst case, captioning spends its
own budget and the document still has most of its own left.

*Model provisioning* -- the stack's existing Ollama, via VISION_MODEL. Empty
(the default) disables the feature entirely: no model pull, no behavior
change. The Compose ollama-model-init only pulls a vision model when
VISION_MODEL is set, and the Helm chart only injects the env var when
`ingestionWorker.visionModel` is set.

*How a caption becomes retrievable* -- one ParsedSection per captioned image,
content_type="image" (issue #89's content-type tagging), page/slide
attribution where the format provides it. Chunking already treats non-"text"
sections as atomic, so a caption is never split, and the content_type lands
in the chunk payload where #89's content_type_boosts can weight it at query
time.
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import io
import logging
import os
from dataclasses import dataclass
from pathlib import Path

import httpx

from app import metrics
from app.parsing import ParsedSection
from common.log_safety import log_safe

logger = logging.getLogger("ingestion-worker")

OLLAMA_URL = os.environ.get("OLLAMA_URL", "http://ollama:11434")
# Empty string = captioning disabled (the default). There is deliberately no
# hardcoded model fallback: pulling a multi-GB vision model must be an
# explicit deployment decision, not a side effect of upgrading the worker.
VISION_MODEL = os.environ.get("VISION_MODEL", "")

MAX_IMAGES_PER_DOCUMENT = int(os.environ.get("MAX_IMAGES_PER_DOCUMENT", "20"))
CAPTIONING_TIMEOUT_SECONDS = float(os.environ.get("CAPTIONING_TIMEOUT_SECONDS", "90"))
CAPTION_REQUEST_TIMEOUT_SECONDS = float(os.environ.get("CAPTION_REQUEST_TIMEOUT_SECONDS", "30"))

# Below either bound an image is almost certainly a bullet glyph, border
# artwork, or a logo -- captioning it costs a model call and adds a chunk of
# retrieval noise ("a blue circle") to the corpus.
MIN_IMAGE_BYTES = int(os.environ.get("CAPTION_MIN_IMAGE_BYTES", "4096"))
MIN_IMAGE_DIMENSION = int(os.environ.get("CAPTION_MIN_IMAGE_DIMENSION", "64"))

# Written for retrieval, not for prettiness: the caption is embedded and
# BM25-indexed verbatim, so it should name the things a person would query
# for (chart values, labels, visible text), not editorialize.
CAPTION_PROMPT = (
    "Describe this image from a document so its content can be found by text "
    "search. State what it shows; transcribe any visible text, labels, and "
    "values; for charts and diagrams, describe the axes, series, and the "
    "relationship shown. Be factual and concise (under 120 words)."
)


@dataclass
class ExtractedImage:
    """Raw image bytes plus the attribution parsing would have given a text
    section from the same place."""

    data: bytes
    page_or_slide: int | None = None


def captioning_enabled() -> bool:
    return bool(VISION_MODEL)


def _extract_pdf_images(content: bytes) -> list[ExtractedImage]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    images: list[ExtractedImage] = []
    for page_number, page in enumerate(reader.pages, start=1):
        # Per-image, not per-document: pypdf decodes each image stream lazily
        # and one undecodable stream (an exotic filter, a truncated object)
        # should cost that image, not the page or the document.
        try:
            page_images = list(page.images)
        except Exception as exc:
            logger.warning(
                "could not enumerate images on PDF page %d: %s: %s",
                page_number,
                type(exc).__name__,
                log_safe(exc),
            )
            continue
        for image_file in page_images:
            try:
                images.append(ExtractedImage(data=image_file.data, page_or_slide=page_number))
            except Exception as exc:
                logger.warning(
                    "could not decode an image on PDF page %d: %s: %s",
                    page_number,
                    type(exc).__name__,
                    log_safe(exc),
                )
    return images


def _extract_docx_images(content: bytes) -> list[ExtractedImage]:
    import docx

    document = docx.Document(io.BytesIO(content))
    images: list[ExtractedImage] = []
    # Relationship walk rather than inline-shape iteration: it sees floating
    # (anchored) pictures too, which `document.inline_shapes` misses. The
    # trade-off is losing document-order position -- but DOCX has no page
    # concept at parse time anyway (parsing.py's sections carry no
    # page_or_slide either), so nothing attributable is being dropped.
    for rel in document.part.rels.values():
        if "image" in rel.reltype and not rel.is_external:
            images.append(ExtractedImage(data=rel.target_part.blob))
    return images


def _extract_pptx_images(content: bytes) -> list[ExtractedImage]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(io.BytesIO(content))
    images: list[ExtractedImage] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    images.append(ExtractedImage(data=shape.image.blob, page_or_slide=slide_number))
                except Exception as exc:
                    logger.warning(
                        "could not read a picture on slide %d: %s: %s",
                        slide_number,
                        type(exc).__name__,
                        log_safe(exc),
                    )
    return images


_EXTRACTORS = {
    ".pdf": _extract_pdf_images,
    ".docx": _extract_docx_images,
    ".pptx": _extract_pptx_images,
}


def _worth_captioning(data: bytes) -> bool:
    """Filter out glyph/logo-sized images and anything Pillow can't decode --
    if Pillow can't read it, the vision model won't either, and the model
    call is the expensive part."""
    if len(data) < MIN_IMAGE_BYTES:
        return False
    try:
        from PIL import Image

        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
    except Exception:
        return False
    return width >= MIN_IMAGE_DIMENSION and height >= MIN_IMAGE_DIMENSION


def extract_images(filename: str, content: bytes) -> list[ExtractedImage]:
    """Deduplicated, size-filtered images from a supported document, capped at
    MAX_IMAGES_PER_DOCUMENT. Returns [] for unsupported formats and for any
    extraction-level failure -- same degrade contract as the rest of this
    module."""
    extractor = _EXTRACTORS.get(Path(filename).suffix.lower())
    if extractor is None:
        return []
    try:
        candidates = extractor(content)
    except Exception as exc:
        # parse_document already succeeded on these bytes, so a hard failure
        # here is an image-layer problem, not a corrupt document -- degrade.
        logger.warning(
            "image extraction failed for a %s document: %s: %s",
            Path(filename).suffix.lower(),
            type(exc).__name__,
            log_safe(exc),
        )
        metrics.images_skipped_total.labels(reason="extraction_error").inc()
        return []

    # Dedupe by content hash: the same logo repeated in every page header or
    # slide master would otherwise be captioned once per occurrence, burning
    # the image budget on identical results.
    seen: set[str] = set()
    images: list[ExtractedImage] = []
    for image in candidates:
        if not _worth_captioning(image.data):
            metrics.images_skipped_total.labels(reason="filtered").inc()
            continue
        digest = hashlib.sha256(image.data).hexdigest()
        if digest in seen:
            metrics.images_skipped_total.labels(reason="duplicate").inc()
            continue
        seen.add(digest)
        if len(images) >= MAX_IMAGES_PER_DOCUMENT:
            metrics.images_skipped_total.labels(reason="over_limit").inc()
            continue
        images.append(image)
    return images


async def _caption_one(client: httpx.AsyncClient, image: ExtractedImage) -> str | None:
    resp = await client.post(
        f"{OLLAMA_URL}/api/generate",
        json={
            "model": VISION_MODEL,
            "prompt": CAPTION_PROMPT,
            "images": [base64.b64encode(image.data).decode("ascii")],
            "stream": False,
        },
    )
    resp.raise_for_status()
    caption = str(resp.json().get("response", "")).strip()
    return caption or None


async def caption_images(filename: str, content: bytes) -> list[ParsedSection]:
    """Extract and caption a document's embedded images, returning one
    content_type="image" ParsedSection per caption. Never raises: any failure
    degrades to returning the captions produced so far (possibly none)."""
    if not captioning_enabled():
        return []

    images = await asyncio.to_thread(extract_images, filename, content)
    if not images:
        return []

    sections: list[ParsedSection] = []
    try:
        async with asyncio.timeout(CAPTIONING_TIMEOUT_SECONDS):
            async with httpx.AsyncClient(timeout=CAPTION_REQUEST_TIMEOUT_SECONDS) as client:
                for image in images:
                    try:
                        caption = await _caption_one(client, image)
                    except (httpx.HTTPError, ValueError, KeyError) as exc:
                        # Model down, model missing, a malformed response --
                        # skip this image, keep trying the rest: a transient
                        # blip shouldn't discard the whole document's figures.
                        logger.warning(
                            "captioning failed for an image (page/slide %s): %s: %s",
                            image.page_or_slide,
                            type(exc).__name__,
                            log_safe(exc),
                        )
                        metrics.images_skipped_total.labels(reason="model_error").inc()
                        continue
                    if caption is None:
                        metrics.images_skipped_total.labels(reason="empty_caption").inc()
                        continue
                    sections.append(
                        ParsedSection(
                            text=caption,
                            page_or_slide=image.page_or_slide,
                            content_type="image",
                        )
                    )
                    metrics.images_captioned_total.inc()
    except TimeoutError:
        # The pass outran its budget -- keep what was captioned, count the
        # rest as skipped so the gap is visible on the dashboard rather than
        # silent (the original sin this issue exists to fix).
        remaining = len(images) - len(sections)
        logger.warning(
            "captioning pass hit its %.0fs budget with %d of %d images done",
            CAPTIONING_TIMEOUT_SECONDS,
            len(sections),
            len(images),
        )
        for _ in range(remaining):
            metrics.images_skipped_total.labels(reason="budget_exhausted").inc()
    except Exception as exc:
        # Belt and braces: the degrade contract ("this module never raises")
        # must hold even for a failure mode not anticipated above.
        logger.warning(
            "captioning pass failed: %s: %s; continuing without remaining captions",
            type(exc).__name__,
            log_safe(exc),
        )
        metrics.images_skipped_total.labels(reason="error").inc()
    return sections
