"""FR-3: parse each supported format into clean text, preserving structural
signal (headings, page/slide numbers) for chunking (FR-4) and citation
(FR-27). FR-9: corrupt/password-protected files raise ParsingError, which the
upload route turns into a clear 4xx instead of a 500 or a silently-empty doc.

Pragmatic choice for this pass: lightweight pure-Python per-format libraries
(pypdf/pdfplumber/python-docx/python-pptx/openpyxl/BeautifulSoup) rather than
the heavier Docling/Unstructured candidates from REQUIREMENTS.md Section 7.4
-- those remain worth evaluating later, but add a large model-download
footprint this dev pass doesn't need.

PDF and DOCX tables are extracted as their own markdown blocks (see
`_table_to_markdown`) rather than left to fall through to plain paragraph/page
text extraction, which flattens a table's rows and columns into an
unstructured word sequence -- e.g. a 3x2 table becomes "Name Role Clearance
Alice Curator Secret" with no indication of which words belonged to which
cell. pdfplumber's table detection is heuristic (it looks for ruling lines /
aligned whitespace) and can occasionally miss a borderless table or, rarely,
false-positive on text that merely looks tabular; this is a known tradeoff of
a lightweight approach, not a guarantee of perfect table recall.
"""

from __future__ import annotations

import io
import logging
import os
import zipfile
from collections.abc import Iterator
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING

from common.log_safety import log_safe

if TYPE_CHECKING:
    from docx.document import Document as DocxDocument
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph
    from pypdf import PdfReader


logger = logging.getLogger("ingestion-worker")


class ParsingError(Exception):
    pass


# NFR-7: .docx/.pptx/.xlsx are ZIP archives under the hood, and none of
# python-docx/python-pptx/openpyxl guard against a zip bomb -- a small,
# maliciously-crafted archive can decompress to gigabytes and exhaust worker
# memory long before MAX_UPLOAD_BYTES (services/ingestion-api/app/routes/
# upload.py) ever sees anything, since that only bounds the *compressed*
# upload size. Checked against the raw zip before handing it to any of those
# libraries.
MAX_ZIP_UNCOMPRESSED_BYTES = 200 * 1024 * 1024  # 200MB decompressed, well over any real OOXML doc
MAX_ZIP_COMPRESSION_RATIO = 200  # legitimate OOXML XML parts rarely exceed the low double digits

# Issue #208: the ZIP guard above covers OOXML only. PDF has its own
# amplification shape -- a small file can declare an enormous number of pages,
# and pdfplumber's per-page work (find_tables in particular, which is
# heuristic and geometric) is expensive enough that page count alone is a
# usable denial-of-service lever well inside MAX_UPLOAD_BYTES.
#
# 5000 pages is far above any real document in this corpus (the largest DoD
# policy PDFs run to the low hundreds) and cheap to check: PdfReader reads the
# page tree without rendering anything, so the count is known before any
# per-page work begins.
MAX_PDF_PAGES = int(os.environ.get("MAX_PDF_PAGES", "5000"))


def _check_zip_bomb(content: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            total_uncompressed = 0
            for info in zf.infolist():
                total_uncompressed += info.file_size
                if info.compress_size and info.file_size / info.compress_size > (
                    MAX_ZIP_COMPRESSION_RATIO
                ):
                    raise ParsingError(
                        f"archive entry '{info.filename}' has a compression ratio consistent "
                        "with a zip bomb"
                    )
                if total_uncompressed > MAX_ZIP_UNCOMPRESSED_BYTES:
                    raise ParsingError(
                        "archive would decompress to over "
                        f"{MAX_ZIP_UNCOMPRESSED_BYTES // (1024 * 1024)}MB, consistent with a "
                        "zip bomb"
                    )
    except zipfile.BadZipFile as exc:
        raise ParsingError(f"corrupt archive: {exc}") from exc


@dataclass
class OcrStatus:
    """Issue #306 gap 3: per-document signal for content that could not be
    scanned for classification markings at all -- OCR unavailable, the
    per-document OCR budget exhausted, or a page's OCR pass recognized no
    text. Without this, those cases and a page with genuinely no markings
    both leave the marking detector with zero findings, and a curator can't
    tell "checked, clean" from "never checked" (see _apply_tagging_advisory
    in app/processing.py, which reads this after parse_document returns)."""

    skipped_pages: int = 0
    reasons: set[str] = field(default_factory=set)

    def record(self, reason: str) -> None:
        self.skipped_pages += 1
        self.reasons.add(reason)

    @property
    def any_skipped(self) -> bool:
        return self.skipped_pages > 0


@dataclass
class ParsedSection:
    """One structural unit of a document -- a heading's worth of text, a PDF
    page, a slide, a spreadsheet sheet. Chunking (FR-4) never splits across
    section boundaries, only within them.

    content_type (issue #89): what kind of content this section holds --
    "text" (prose, the default), "table" (a markdown table block extracted
    by `_table_to_markdown`, or a spreadsheet sheet, which is tabular by
    construction), or "image" (a VLM-generated caption of an embedded figure,
    issue #92 -- produced by app/captioning.py after parsing, not by this
    module). Tables are emitted as their own section rather than joined
    into the surrounding prose specifically so chunking never mixes the two
    within one chunk -- that's what lets each chunk carry a single, accurate
    content_type through to the Qdrant payload instead of a per-chunk guess."""

    text: str
    heading: str | None = None
    page_or_slide: int | None = None
    content_type: str = "text"


def parse_document(
    filename: str, content: bytes, ocr_status: OcrStatus | None = None
) -> list[ParsedSection]:
    """`ocr_status`, when passed, is mutated in place with any scanned-page
    OCR fallback (issue #241) content that could not be recovered (issue
    #306 gap 3) -- see OcrStatus above. Callers that don't care (most tests,
    every non-PDF format) can omit it."""
    if not content:
        raise ParsingError("empty file")

    ext = Path(filename).suffix.lower()
    try:
        if ext in (".txt",):
            return _parse_txt(content)
        if ext in (".md", ".markdown"):
            return _parse_markdown(content)
        if ext in (".html", ".htm"):
            return _parse_html(content)
        if ext == ".pdf":
            return _parse_pdf(content, ocr_status)
        if ext == ".docx":
            _check_zip_bomb(content)
            return _parse_docx(content)
        if ext == ".pptx":
            _check_zip_bomb(content)
            return _parse_pptx(content)
        if ext == ".xlsx":
            _check_zip_bomb(content)
            return _parse_xlsx(content)
        if ext in (".png", ".jpg", ".jpeg", ".tif", ".tiff"):
            return _parse_image(content)
    except ParsingError:
        # Raised deliberately by this module -- "password-protected PDF",
        # "unsupported file type", the page/zip bounds. Those messages are
        # written for the uploader and are exactly what FR-8 should return.
        raise
    except Exception as exc:
        # #214: anything else is a parser library's own exception, and its
        # text carries library internals, file offsets, and sometimes
        # fragments of the document. It is persisted to
        # Document.processing_error and re-served through the status endpoint,
        # so it reaches whoever uploaded the file.
        #
        # The uploader needs to know their document could not be read and
        # that re-uploading the same bytes will not help; they do not need
        # pypdf's stack vocabulary. The detail goes to the log, where the
        # operator who can act on it is.
        logger.warning(
            "parser failed on a %s document: %s: %s",
            ext or "unknown-type",
            type(exc).__name__,
            log_safe(exc),
        )
        raise ParsingError(
            f"could not read this {ext or 'file'} -- it may be corrupt or in an unexpected format"
        ) from exc

    raise ParsingError(f"unsupported file type: {ext or filename}")


def _parse_txt(content: bytes) -> list[ParsedSection]:
    text = content.decode("utf-8", errors="replace")
    return [ParsedSection(text=text)]


def _parse_markdown(content: bytes) -> list[ParsedSection]:
    text = content.decode("utf-8", errors="replace")
    sections: list[ParsedSection] = []
    current_heading: str | None = None
    current_lines: list[str] = []

    def flush() -> None:
        body = "\n".join(current_lines).strip()
        if body:
            sections.append(ParsedSection(text=body, heading=current_heading))

    for line in text.splitlines():
        if line.lstrip().startswith("#"):
            flush()
            current_heading = line.lstrip("#").strip() or None
            current_lines = []
        else:
            current_lines.append(line)
    flush()

    return sections or [ParsedSection(text=text)]


def _parse_html(content: bytes) -> list[ParsedSection]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()

    headings = soup.find_all(["h1", "h2", "h3"])
    if not headings:
        text = soup.get_text(separator="\n", strip=True)
        return [ParsedSection(text=text)] if text else []

    sections: list[ParsedSection] = []
    for heading in headings:
        heading_text = heading.get_text(strip=True)
        body_parts = []
        for sibling in heading.find_next_siblings():
            if getattr(sibling, "name", None) in ("h1", "h2", "h3"):
                break
            body_parts.append(sibling.get_text(separator="\n", strip=True))
        body = "\n".join(p for p in body_parts if p)
        if body:
            sections.append(ParsedSection(text=body, heading=heading_text))
    return sections


def _parse_image(content: bytes) -> list[ParsedSection]:
    """Issue #241: a standalone image upload -- a scanned memo, a photographed
    whiteboard. OCR is the only content path for this format, so unlike the
    scanned-PDF fallback below, failure here is an FR-8 error the uploader can
    act on, not a silent skip: an empty "success" would be a lie."""
    from app import ocr

    if not ocr.ocr_available():
        raise ParsingError(
            "this deployment cannot read image files (OCR is unavailable); "
            "supply the content as a text-bearing format instead"
        )
    text = ocr.ocr_image(content)
    if not text:
        raise ParsingError(
            "no readable text was recognized in this image -- if it holds a "
            "figure rather than text, embed it in a document instead"
        )
    # content_type "ocr" (issue #89's tagging): provenance for curators
    # reviewing machine-read text, and weightable via CONTENT_TYPE_BOOSTS.
    # Chunked by the normal sliding window (see chunking.py) -- OCR output is
    # prose, unlike the atomic table/image section kinds.
    return [ParsedSection(text=text, content_type="ocr")]


def _table_to_markdown(grid: list[list[str | None]]) -> str:
    """Render a pdfplumber/python-docx table grid (rows of cell strings, may
    contain None for empty cells) as a markdown table. Rows that are entirely
    empty (a common pdfplumber artifact around merged/spanning cells) are
    dropped rather than emitted as blank table rows."""
    rows = [[("" if cell is None else str(cell).strip()) for cell in row] for row in grid]
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return ""

    width = len(rows[0])
    lines = [
        "| " + " | ".join(rows[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _parse_pdf(content: bytes, ocr_status: OcrStatus | None = None) -> list[ParsedSection]:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(io.BytesIO(content))
    except PdfReadError as exc:
        raise ParsingError(f"corrupt PDF: {exc}") from exc

    # Try an empty password (common for "restricted" rather than truly
    # encrypted PDFs); anything else is reported as password-protected.
    if reader.is_encrypted and reader.decrypt("") == 0:
        raise ParsingError("password-protected PDF")

    # Checked before pdfplumber opens the file, so a page-count bomb is
    # rejected without doing any per-page geometry work. This bounds one
    # amplification shape, not all of them -- a single page can still hold
    # pathological content, which is what the worker's
    # PROCESSING_TIMEOUT_SECONDS budget is for. The two are complementary:
    # this one gives a precise, actionable FR-8 message, the timeout is the
    # catch-all.
    page_count = len(reader.pages)
    if page_count > MAX_PDF_PAGES:
        raise ParsingError(f"PDF has {page_count} pages, over the {MAX_PDF_PAGES}-page limit")

    import pdfplumber

    sections = []
    # Issue #241: shared OCR budget for this document's scanned-page fallback,
    # decremented per image OCR'd (see ocr.py for why it exists).
    ocr_budget = _MAX_OCR_IMAGES
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for i, page in enumerate(pdf.pages):
            tables = page.find_tables()

            # Extract prose from everything outside the detected table
            # bounding boxes, so table cell text isn't also duplicated (and
            # scrambled) into the surrounding paragraph text.
            prose_page = page
            for table in tables:
                prose_page = prose_page.outside_bbox(table.bbox)
            prose = (prose_page.extract_text() or "").strip()

            # issue #89: prose and each table become their own section
            # (same page_or_slide, content_type "text" vs "table") instead of
            # being joined into one blob -- see the ParsedSection docstring.
            if prose:
                sections.append(ParsedSection(text=prose, page_or_slide=i + 1, content_type="text"))
            for table in tables:
                markdown = _table_to_markdown(table.extract())
                if markdown:
                    sections.append(
                        ParsedSection(text=markdown, page_or_slide=i + 1, content_type="table")
                    )

            # Issue #241: scanned-page fallback. Fires only when this page
            # contributed nothing at all -- no prose, no tables -- so a PDF
            # with a text layer parses byte-identically to before. A page
            # that *has* text plus a figure is #92 captioning's territory,
            # not OCR's.
            if not prose and not tables:
                ocr_sections, ocr_budget = _ocr_pdf_page(reader, i, ocr_budget, ocr_status)
                sections.extend(ocr_sections)

    return sections


# Read once at import (matches ocr.py's own constant) so the per-document
# budget in _parse_pdf is testable by monkeypatching this module.
def _default_ocr_budget() -> int:
    from app.ocr import MAX_OCR_IMAGES_PER_DOCUMENT

    return MAX_OCR_IMAGES_PER_DOCUMENT


_MAX_OCR_IMAGES = _default_ocr_budget()


def _ocr_pdf_page(
    reader: PdfReader, page_index: int, budget: int, ocr_status: OcrStatus | None = None
) -> tuple[list[ParsedSection], int]:
    """OCR the embedded images of one text-less PDF page. Degrade-only: a
    missing tesseract, an undecodable image, or an exhausted budget costs
    text that today contributes nothing anyway -- logged, never raised.

    `ocr_status`, when passed, records *why* a page ended up contributing no
    text (issue #306 gap 3) -- the marking-mismatch advisory reads this to
    tell a curator "this content was never scanned" apart from "it was
    scanned and carries no markings"."""
    from app import ocr

    if not ocr.ocr_available():
        # One line per page, so a fully scanned document shows the gap
        # loudly in the logs rather than failing with a bare "no
        # extractable text" and no explanation.
        logger.warning(
            "PDF page %d has no text layer and OCR is unavailable; its content is skipped",
            page_index + 1,
        )
        if ocr_status is not None:
            ocr_status.record("ocr_unavailable")
        return [], budget
    try:
        page_images = list(reader.pages[page_index].images)
    except Exception as exc:
        logger.warning(
            "could not read images on text-less PDF page %d: %s: %s",
            page_index + 1,
            type(exc).__name__,
            log_safe(exc),
        )
        if ocr_status is not None:
            ocr_status.record("unreadable_images")
        return [], budget

    texts: list[str] = []
    budget_exhausted = False
    for image_file in page_images:
        if budget <= 0:
            logger.warning(
                "OCR image budget exhausted at PDF page %d; remaining scanned pages are skipped",
                page_index + 1,
            )
            budget_exhausted = True
            break
        budget -= 1
        try:
            text = ocr.ocr_image(image_file.data)
        except Exception as exc:  # ocr_image shouldn't raise; belt and braces
            logger.warning(
                "OCR failed on PDF page %d: %s: %s",
                page_index + 1,
                type(exc).__name__,
                log_safe(exc),
            )
            continue
        if text:
            texts.append(text)
    if not texts:
        if ocr_status is not None:
            ocr_status.record("budget_exhausted" if budget_exhausted else "no_text_recognized")
        return [], budget
    return [
        ParsedSection(text="\n".join(texts), page_or_slide=page_index + 1, content_type="ocr")
    ], budget


def _iter_docx_block_items(document: DocxDocument) -> Iterator[DocxParagraph | DocxTable]:
    """Yield a document's paragraphs and tables in document order.

    python-docx's own `document.paragraphs` and `document.tables` are two
    separate flat lists with no shared ordering, so iterating them
    independently loses each table's position relative to the surrounding
    paragraphs. Walking `document.element.body`'s direct children and
    wrapping each one back into its python-docx object is the standard
    workaround for recovering document order.
    """
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            yield Paragraph(child, document)
        elif child.tag == qn("w:tbl"):
            yield Table(child, document)


def _docx_table_grid(table: DocxTable) -> list[list[str | None]]:
    # Note: python-docx repeats the same Cell object across a merged span, so
    # a horizontally- or vertically-merged cell's text appears once per
    # spanned position rather than once per visual cell -- a known
    # python-docx quirk, not something this function corrects for.
    return [[cell.text for cell in row.cells] for row in table.rows]


def _parse_docx(content: bytes) -> list[ParsedSection]:
    import docx
    from docx.table import Table

    document = docx.Document(io.BytesIO(content))
    sections: list[ParsedSection] = []
    current_heading: str | None = None
    current_lines: list[str] = []

    def flush_prose() -> None:
        body = "\n".join(current_lines).strip()
        if body:
            sections.append(ParsedSection(text=body, heading=current_heading, content_type="text"))
        current_lines.clear()

    for block in _iter_docx_block_items(document):
        if isinstance(block, Table):
            markdown = _table_to_markdown(_docx_table_grid(block))
            if markdown:
                # issue #89: flush any prose accumulated so far as its own
                # "text" section, then the table as its own "table" section --
                # prose before/after a table under the same heading no longer
                # get glued to it into one section (see ParsedSection
                # docstring for why that matters for chunk-level tagging).
                flush_prose()
                sections.append(
                    ParsedSection(text=markdown, heading=current_heading, content_type="table")
                )
        elif block.style and block.style.name.startswith("Heading"):
            flush_prose()
            current_heading = block.text.strip() or None
        elif block.text.strip():
            current_lines.append(block.text)
    flush_prose()

    return sections


def _parse_pptx(content: bytes) -> list[ParsedSection]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    sections = []
    for i, slide in enumerate(presentation.slides):
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape and title_shape.text else None
        texts = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue  # already captured as the section heading
            if shape.has_text_frame and shape.text.strip():
                texts.append(shape.text.strip())
        # Fall back to the title as the body when a slide has no other text
        # (e.g. a section-divider slide) so it isn't dropped from the corpus.
        body = "\n".join(texts) or title or ""
        if body:
            sections.append(ParsedSection(text=body, heading=title, page_or_slide=i + 1))
    return sections


def _parse_xlsx(content: bytes) -> list[ParsedSection]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    sections = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        body = "\n".join(rows)
        if body:
            # issue #89: a spreadsheet sheet is tabular by construction, no
            # detection heuristic needed.
            sections.append(ParsedSection(text=body, heading=sheet.title, content_type="table"))
    return sections
