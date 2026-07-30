"""Issue #92: VLM image captioning -- extraction, filtering, and the
degrade-on-failure contract (caption_images never raises; a captioning
failure costs captions, never the document).

No live Ollama anywhere: the model endpoint is mocked with respx. Fixture
documents are authored in memory with Pillow/python-docx/python-pptx, the
same author-the-input approach the parsing tests use.
"""

from __future__ import annotations

import io
import json

import httpx
import pytest
import respx

from app import captioning
from app.captioning import caption_images, extract_images

OLLAMA = "http://ollama:11434"


# --- fixture authoring --------------------------------------------------------


def _png_bytes(width: int = 256, height: int = 160, seed: int = 0) -> bytes:
    """A deterministic noise image: solid colors compress under
    MIN_IMAGE_BYTES and would be filtered as glyph-sized; noise doesn't.
    Different seeds give different bytes (for the dedupe test)."""
    import random

    from PIL import Image

    rng = random.Random(seed)
    image = Image.new("RGB", (width, height))
    image.putdata(
        [
            (rng.randrange(256), rng.randrange(256), rng.randrange(256))
            for _ in range(width * height)
        ]
    )
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    return buf.getvalue()


def _pdf_with_images(image_pages: list[bytes]) -> bytes:
    """One page per image, authored by Pillow's PDF writer -- the smallest way
    to get a real PDF with real embedded image XObjects."""
    from PIL import Image

    pages = [Image.open(io.BytesIO(data)).convert("RGB") for data in image_pages]
    buf = io.BytesIO()
    pages[0].save(buf, format="PDF", save_all=True, append_images=pages[1:])
    return buf.getvalue()


def _docx_with_images(images: list[bytes]) -> bytes:
    import docx

    document = docx.Document()
    document.add_paragraph("Prose before the figure.")
    for data in images:
        document.add_picture(io.BytesIO(data))
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _pptx_with_images(images_per_slide: list[list[bytes]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for slide_images in images_per_slide:
        slide = presentation.slides.add_slide(blank)
        for data in slide_images:
            slide.shapes.add_picture(io.BytesIO(data), Inches(1), Inches(1))
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


@pytest.fixture
def enabled(monkeypatch):
    monkeypatch.setattr(captioning, "VISION_MODEL", "test-vision-model")
    monkeypatch.setattr(captioning, "OLLAMA_URL", OLLAMA)


def _mock_caption(text: str = "A bar chart of quarterly totals."):
    return respx.post(f"{OLLAMA}/api/generate").mock(
        return_value=httpx.Response(200, json={"response": text})
    )


# --- extraction ---------------------------------------------------------------


def test_extract_pdf_images_with_page_attribution():
    first, second = _png_bytes(seed=1), _png_bytes(seed=2)
    images = extract_images("doc.pdf", _pdf_with_images([first, second]))
    assert [image.page_or_slide for image in images] == [1, 2]


def test_extract_docx_images():
    images = extract_images("doc.docx", _docx_with_images([_png_bytes()]))
    assert len(images) == 1
    # DOCX has no page concept at parse time (parsing.py's sections carry
    # no page_or_slide either), so attribution is deliberately absent.
    assert images[0].page_or_slide is None


def test_extract_pptx_images_with_slide_attribution():
    content = _pptx_with_images([[_png_bytes(seed=1)], [], [_png_bytes(seed=2)]])
    images = extract_images("deck.pptx", content)
    assert [image.page_or_slide for image in images] == [1, 3]


def test_unsupported_formats_yield_nothing():
    assert extract_images("notes.txt", b"plain text, no images") == []
    assert extract_images("notes.md", b"# heading") == []


def test_tiny_images_filtered():
    # 16x16 is under MIN_IMAGE_DIMENSION -- a bullet glyph, not a figure.
    content = _docx_with_images([_png_bytes(16, 16)])
    assert extract_images("doc.docx", content) == []


def test_duplicate_images_extracted_once():
    same = _png_bytes()
    content = _pptx_with_images([[same], [same], [same]])
    assert len(extract_images("deck.pptx", content)) == 1


def test_image_count_capped(monkeypatch):
    monkeypatch.setattr(captioning, "MAX_IMAGES_PER_DOCUMENT", 2)
    seeds = [1, 2, 3, 4]
    content = _pptx_with_images([[_png_bytes(seed=s)] for s in seeds])
    assert len(extract_images("deck.pptx", content)) == 2


def test_extraction_failure_degrades_to_empty():
    # A .docx that is not a zip at all: python-docx raises, extract_images
    # must swallow it (parse-level rejection is parsing.py's job, not ours).
    assert extract_images("doc.docx", b"not a zip archive") == []


# --- captioning ---------------------------------------------------------------


async def test_disabled_is_a_no_op_without_any_call():
    assert captioning.captioning_enabled() is False
    with respx.mock(assert_all_called=False) as router:
        route = router.post(f"{OLLAMA}/api/generate")
        assert await caption_images("doc.pdf", _pdf_with_images([_png_bytes()])) == []
        assert not route.called


@respx.mock
async def test_captions_become_image_sections(enabled):
    route = _mock_caption("A line chart: throughput rising from 10 to 90.")
    sections = await caption_images("doc.pdf", _pdf_with_images([_png_bytes()]))
    assert route.called
    assert len(sections) == 1
    assert sections[0].content_type == "image"
    assert sections[0].page_or_slide == 1
    assert "throughput rising" in sections[0].text


@respx.mock
async def test_request_carries_model_prompt_and_image(enabled):
    route = _mock_caption()
    await caption_images("doc.pdf", _pdf_with_images([_png_bytes()]))
    body = json.loads(route.calls.last.request.content)
    assert body["model"] == "test-vision-model"
    assert body["prompt"] == captioning.CAPTION_PROMPT
    assert body["stream"] is False
    assert len(body["images"]) == 1 and isinstance(body["images"][0], str)


@respx.mock
async def test_model_unreachable_degrades_to_no_captions(enabled):
    respx.post(f"{OLLAMA}/api/generate").mock(side_effect=httpx.ConnectError("refused"))
    assert await caption_images("doc.pdf", _pdf_with_images([_png_bytes()])) == []


@respx.mock
async def test_one_failing_image_does_not_lose_the_rest(enabled):
    responses = iter(
        [
            httpx.Response(500, text="model exploded"),
            httpx.Response(200, json={"response": "A network topology diagram."}),
        ]
    )
    respx.post(f"{OLLAMA}/api/generate").mock(side_effect=lambda request: next(responses))
    content = _pdf_with_images([_png_bytes(seed=1), _png_bytes(seed=2)])
    sections = await caption_images("doc.pdf", content)
    assert [s.text for s in sections] == ["A network topology diagram."]
    assert sections[0].page_or_slide == 2


@respx.mock
async def test_empty_caption_is_skipped(enabled):
    _mock_caption("   ")
    assert await caption_images("doc.pdf", _pdf_with_images([_png_bytes()])) == []


@respx.mock
async def test_budget_exhaustion_keeps_partial_results(enabled, monkeypatch):
    # First caption lands, then the pass's overall budget expires: the
    # produced caption must survive, and nothing may raise.
    monkeypatch.setattr(captioning, "CAPTIONING_TIMEOUT_SECONDS", 0.25)

    calls = 0

    async def slow_then_slower(request):
        nonlocal calls
        calls += 1
        if calls > 1:
            import asyncio

            await asyncio.sleep(1.0)
        return httpx.Response(200, json={"response": f"Caption {calls}."})

    respx.post(f"{OLLAMA}/api/generate").mock(side_effect=slow_then_slower)
    content = _pdf_with_images([_png_bytes(seed=1), _png_bytes(seed=2)])
    sections = await caption_images("doc.pdf", content)
    assert [s.text for s in sections] == ["Caption 1."]
