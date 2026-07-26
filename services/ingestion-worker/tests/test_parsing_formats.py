"""Unit tests for the ingestion-worker's parsing (FR-3/FR-9/NFR-7): text and
markdown paths, HTML script stripping, corrupt/unsupported input errors, and
the zip-bomb guard in front of the OOXML parsers.
"""

from __future__ import annotations

import io
import zipfile

import pytest

from app.parsing import ParsingError, _check_zip_bomb, parse_document


class TestTextParsing:
    def test_txt_round_trip(self):
        sections = parse_document("notes.txt", b"hello world")
        assert len(sections) == 1
        assert sections[0].text == "hello world"
        assert sections[0].heading is None

    def test_invalid_utf8_replaced_not_raised(self):
        sections = parse_document("notes.txt", b"ok \xff\xfe bytes")
        assert "ok" in sections[0].text

    def test_empty_file_rejected(self):
        with pytest.raises(ParsingError, match="empty file"):
            parse_document("notes.txt", b"")


class TestMarkdownParsing:
    def test_headings_split_sections(self):
        content = b"# Title\n\nintro text\n\n## Details\n\nbody text\n"
        sections = parse_document("doc.md", content)
        assert [(s.heading, s.text) for s in sections] == [
            ("Title", "intro text"),
            ("Details", "body text"),
        ]

    def test_headingless_markdown_single_section(self):
        sections = parse_document("doc.md", b"just some text")
        assert len(sections) == 1
        assert sections[0].heading is None

    def test_leading_content_before_first_heading_dropped_from_sections(self):
        # Content before the first heading is flushed only if non-empty; it is
        # kept as an untitled section.
        sections = parse_document("doc.md", b"preamble\n\n# H1\nbody\n")
        assert sections[0].heading is None
        assert "preamble" in sections[0].text
        assert sections[1].heading == "H1"


class TestHtmlParsing:
    def test_scripts_and_styles_stripped(self):
        html = b"<html><head><style>body{}</style></head><body>" \
               b"<p>visible</p><script>alert(1)</script></body></html>"
        sections = parse_document("page.html", html)
        text = " ".join(s.text for s in sections)
        assert "visible" in text
        assert "alert" not in text
        assert "body{}" not in text

    def test_headings_define_sections(self):
        html = b"<h1>One</h1><p>first</p><h2>Two</h2><p>second</p>"
        sections = parse_document("page.html", html)
        assert [(s.heading, s.text) for s in sections] == [
            ("One", "first"),
            ("Two", "second"),
        ]


class TestUnsupportedAndCorrupt:
    def test_unsupported_extension_rejected(self):
        with pytest.raises(ParsingError, match="unsupported file type"):
            parse_document("archive.tar.gz", b"data")

    def test_corrupt_pdf_rejected(self):
        with pytest.raises(ParsingError):
            parse_document("broken.pdf", b"this is not a pdf")

    def test_corrupt_docx_rejected_as_archive(self):
        with pytest.raises(ParsingError, match="corrupt archive"):
            parse_document("broken.docx", b"definitely not a zip")


class TestOfficeFormats:
    def test_docx_headings_and_paragraphs(self):
        import docx

        document = docx.Document()
        document.add_heading("Overview", level=1)
        document.add_paragraph("first body paragraph")
        document.add_heading("Details", level=2)
        document.add_paragraph("second body paragraph")
        buf = io.BytesIO()
        document.save(buf)

        sections = parse_document("doc.docx", buf.getvalue())
        assert [(s.heading, s.text) for s in sections] == [
            ("Overview", "first body paragraph"),
            ("Details", "second body paragraph"),
        ]

    def test_pptx_slides_with_titles(self):
        from pptx import Presentation

        presentation = Presentation()
        layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = "Slide One"
        slide.placeholders[1].text = "slide one body"
        buf = io.BytesIO()
        presentation.save(buf)

        sections = parse_document("slides.pptx", buf.getvalue())
        assert len(sections) == 1
        assert sections[0].heading == "Slide One"
        assert "slide one body" in sections[0].text
        assert sections[0].page_or_slide == 1

    def test_xlsx_rows_become_pipe_separated_text(self):
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Inventory"
        sheet.append(["item", "count"])
        sheet.append(["rifles", 10])
        buf = io.BytesIO()
        workbook.save(buf)

        sections = parse_document("sheet.xlsx", buf.getvalue())
        assert len(sections) == 1
        assert sections[0].heading == "Inventory"
        assert "item | count" in sections[0].text
        assert "rifles | 10" in sections[0].text


class TestZipBombGuard:
    def _zip_with(self, payload: bytes, name: str = "word/document.xml") -> bytes:
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr(name, payload)
        return buf.getvalue()

    def test_extreme_compression_ratio_rejected(self):
        # 10MB of zeros compresses to nearly nothing -- far past the ratio cap.
        archive = self._zip_with(b"\x00" * (10 * 1024 * 1024))
        with pytest.raises(ParsingError, match="zip bomb"):
            _check_zip_bomb(archive)

    def test_legitimate_archive_passes(self):
        archive = self._zip_with(b'<w:document>real content</w:document>' * 50)
        _check_zip_bomb(archive)  # no exception

    def test_non_zip_rejected(self):
        with pytest.raises(ParsingError, match="corrupt archive"):
            _check_zip_bomb(b"not a zip at all")

    def test_guard_runs_before_docx_parse(self):
        archive = self._zip_with(b"\x00" * (10 * 1024 * 1024))
        with pytest.raises(ParsingError, match="zip bomb"):
            parse_document("evil.docx", archive)
