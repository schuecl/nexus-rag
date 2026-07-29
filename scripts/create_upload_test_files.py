"""Create a harmless multi-format corpus for exercising document uploads.

Run from the repository root after installing the ingestion-worker's parser
dependencies.  The generated files contain memorable fictional facts so a
human can ask clear retrieval questions after they are approved.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = REPOSITORY_ROOT / "sample-data" / "upload-tests"


def _pdf_escape(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def write_pdf(path: Path) -> None:
    """Write a small standards-compliant PDF without adding a PDF dependency."""
    lines = [
        "Blue Lantern Exercise Brief",
        "The fictional Blue Lantern readiness drill occurs on the second Tuesday",
        "of every month at 09:30 in Building 42.",
        "Participants bring a badge, notebook, and completed safety checklist.",
        "The exercise coordinator is the Training Office.",
    ]
    commands = ["BT", "/F1 17 Tf", "72 740 Td"]
    for index, line in enumerate(lines):
        if index:
            commands.extend(["0 -28 Td", "/F1 12 Tf"])
        commands.append(f"({_pdf_escape(line)}) Tj")
    commands.append("ET")
    stream = "\n".join(commands).encode("ascii")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>"
        ),
        (
            b"<< /Length "
            + str(len(stream)).encode("ascii")
            + b" >>\nstream\n"
            + stream
            + b"\nendstream"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    output = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, obj in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{number} 0 obj\n".encode("ascii"))
        output.extend(obj)
        output.extend(b"\nendobj\n")
    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    path.write_bytes(output)


def write_docx(path: Path) -> None:
    document = Document()
    document.add_heading("Equipment Checkout Guide", level=1)
    document.add_paragraph(
        "Training laptops may be checked out for a maximum of 14 calendar days. "
        "Extensions require approval from the Logistics Desk."
    )
    document.add_heading("Return checklist", level=2)
    for item in ("Return the charger", "Remove personal files", "Report visible damage"):
        document.add_paragraph(item, style="List Bullet")
    document.save(path)


def write_pptx(path: Path) -> None:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Project Northwind Test Schedule"
    title_slide.placeholders[1].text = "Fictional milestones for retrieval testing"

    schedule_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    schedule_slide.shapes.title.text = "Milestones"
    schedule_slide.placeholders[
        1
    ].text = "Alpha review: August 15\nBeta review: September 30\nTraining release: October 20"
    presentation.save(path)


def write_xlsx(path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Test Inventory"
    sheet.append(["Item", "Location", "Quantity", "Inspection cycle"])
    sheet.append(["Orion radio", "Warehouse A", 18, "Monthly"])
    sheet.append(["Cobalt tablet", "Warehouse B", 24, "Quarterly"])
    sheet.append(["Amber projector", "Training Room 3", 6, "Semiannual"])
    workbook.save(path)


def write_text_formats(output_dir: Path) -> None:
    (output_dir / "support-hours.txt").write_text(
        "Nexus test support hours\n\n"
        "The fictional support desk is open Monday through Thursday from "
        "07:30 to 16:30 and Friday from 07:30 to 13:00.\n",
        encoding="utf-8",
    )
    (output_dir / "test-glossary.md").write_text(
        "# Test Corpus Glossary\n\n"
        "- **Blue Lantern**: the monthly readiness drill held in Building 42.\n"
        "- **Northwind**: the fictional training release scheduled for October 20.\n"
        "- **Orion radio**: an inventory item stored in Warehouse A.\n",
        encoding="utf-8",
    )
    (output_dir / "onboarding-checklist.html").write_text(
        '<!doctype html><html lang="en"><head><meta charset="utf-8">'
        "<title>Test Onboarding Checklist</title></head><body>"
        "<h1>Test Onboarding Checklist</h1>"
        "<p>New test participants complete orientation in Room 214.</p>"
        "<ol><li>Collect a visitor badge.</li><li>Read the safety card.</li>"
        "<li>Meet the test coordinator at 10:15.</li></ol></body></html>",
        encoding="utf-8",
    )


def write_images(source_png: Path, output_dir: Path) -> None:
    png_target = output_dir / "unsupported-image-sample.png"
    with Image.open(source_png) as source:
        source.convert("RGB").save(png_target, "PNG", optimize=True)
        source.convert("RGB").save(
            output_dir / "unsupported-image-sample.jpg",
            "JPEG",
            quality=88,
            optimize=True,
        )


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--source-image",
        type=Path,
        required=True,
        help="PNG image used to create the PNG/JPG unsupported-format samples",
    )
    args = parser.parse_args()

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    write_pdf(OUTPUT_DIR / "blue-lantern-brief.pdf")
    write_docx(OUTPUT_DIR / "equipment-checkout-guide.docx")
    write_pptx(OUTPUT_DIR / "northwind-schedule.pptx")
    write_xlsx(OUTPUT_DIR / "test-inventory.xlsx")
    write_text_formats(OUTPUT_DIR)
    write_images(args.source_image.resolve(), OUTPUT_DIR)
    print(f"Created upload test files in {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
